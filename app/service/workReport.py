import copy
import os
from app.core.db.base import Session
from app.util.redis import RedisClient
from app.util.smtpMessage import SendMessage
from app.util.webex import Messages
from app.core.db.models.users import Users
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.autoshape import Shape
from app.util.createDate import CreateDatetime

CONN = RedisClient(1)
SERVICE_NAME = "WorkReport"
DEFALUT_FROM = "./workReport/form.pptx"


class WorkReport:
    async def service(message, roomId, conn):
        redis_value = CONN.getContent(roomId)

        content = Messages().downloadFile(message["files"][0])
        Messages.saveFiletoPptx(content, redis_value["user_name"])
        value = Session
        return conn.postMessage(
            roomId,
            "</br> <h4> 업무보고 파일을 정상적으로 서버에 저장하였어요! \n 수정하시고 싶으시다면 '업무보고' 서비스를 통해 다시 파일을 전송해주세요! \n 종합 마감은 금요일 10시입니다!",
        )

    async def returnMessage(value):
        return "</br> <h4> 업무보고용 파일을 첨부해주세요! 파일형식은 pptx만 가능하답니다"

    def checkValue(message):
        if "files" not in message:
            return False
        else:
            return True


class MergeWorkReport:
    def makeMergePPTX():
        datevlaue = CreateDatetime.today()
        weekValue = CreateDatetime.todayToWeek(
            datevlaue["year"], datevlaue["month"], datevlaue["day"]
        )
        formPptObj = Presentation(DEFALUT_FROM)
        fileList = os.listdir(f"./workReport/{weekValue}주차")

        for cnt, file in enumerate(fileList):
            copyList = MergeWorkReport.copySlide(file, weekValue)
            value = formPptObj.slides[2 + cnt]
            for shape in copyList:
                el = shape
                newel = copy.deepcopy(el)
                value.shapes._spTree.insert_element_before(newel, "p:extLst")

            formPptObj.save(
                "./workReport/{}.{}.{}주차_업무보고_종합.pptx".format(
                    datevlaue["year"], datevlaue["month"], weekValue
                )
            )

    def copySlide(path, week):
        pptObj = Presentation(f"./workReport/{week}주차/{path}")
        copyValue = pptObj.slides[2]
        return_value = []
        for shape in copyValue.shapes:
            if type(shape) == Shape or type(shape) == GraphicFrame:
                return_value.append(copy.deepcopy(shape.element))
        return return_value

    def sendMail(teamName):
        datevlaue = CreateDatetime.today()
        weekValue = CreateDatetime.todayToWeek(
            datevlaue["year"], datevlaue["month"], datevlaue["day"]
        )
        input_value = {
            "title": "{}.{}.{} {} 주간업무 보고 종합입니다.".format(
                datevlaue["year"], datevlaue["month"], datevlaue["day"], teamName
            ),
            "content": """
            
                안녕하십니까, {} 최충은 프로입니다.

                {}년 {}월 {}주차 주간업무보고서 종합해서 보내드립니다.

                감사합니다
            """.format(
                teamName, datevlaue["year"], datevlaue["month"], weekValue
            ),
            "filePath": "./workReport/{}.{}.{}주차_업무보고_종합.pptx".format(
                datevlaue["year"], datevlaue["month"], weekValue
            ),
        }
        SendMessage().writeEmail(**input_value)
